import uvicorn
from fastapi import FastAPI, HTTPException, UploadFile, File, Form
from typing import Optional
from pathlib import Path
from pydantic import BaseModel
from fastapi.responses import FileResponse
from docx import Document
from pptx import Presentation
from bs4 import BeautifulSoup
import pandas as pd
import pdfplumber
import subprocess
from xhtml2pdf import pisa

app = FastAPI()

# Actualizamos la carpeta donde se almacenan los archivos a la ruta del volumen compartido
FILES_FOLDER = Path("/app/shared_files/")

class FileContentResponse(BaseModel):
    """Response model for extracted file text."""
    filename: str
    file_content: str
    additional_text: Optional[str] = None

def read_file_text(file_path: Path) -> str:
    """Extracts the plain text from a file."""
    file_extension = file_path.suffix.lower()
    try:
        if file_extension == ".txt":
            return file_path.read_text(encoding="utf-8")
        elif file_extension == ".pdf":
            with pdfplumber.open(file_path) as pdf:
                return "\n".join(page.extract_text() for page in pdf.pages if page.extract_text())
        elif file_extension == ".docx":
            return "\n".join(p.text for p in Document(file_path).paragraphs)
        elif file_extension == ".pptx":
            return "\n".join(shape.text for slide in Presentation(file_path).slides for shape in slide.shapes if hasattr(shape, "text"))
        elif file_extension == ".html":
            with open(file_path, "r", encoding="utf-8") as f:
                return BeautifulSoup(f, "html.parser").get_text()
        elif file_extension in [".xls", ".xlsx"]:
            excel_data = pd.read_excel(file_path, sheet_name=None)
            return "\n\n".join(f"Sheet: {sheet_name}\n{df.to_string(index=False)}" for sheet_name, df in excel_data.items())
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Error processing file: {str(e)}")
    raise HTTPException(status_code=400, detail="Unsupported file format.")

@app.post("/files/upload/", 
          response_model=FileContentResponse, 
          summary="Upload and extract file content", 
          description="Uploads a file and extracts its text.")
async def upload_file(file: UploadFile = File(...), additional_text: Optional[str] = Form(None)):
    file_path = FILES_FOLDER / file.filename 
    try:
        with open(file_path, "wb") as f:
            f.write(await file.read())
        file_text = read_file_text(file_path)
        return FileContentResponse(filename=file.filename, file_content=file_text, additional_text=additional_text)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing the file: {str(e)}")


@app.get(
    "/files/download/{filename}",
    response_class=FileResponse,
    summary="Download a file",
    description="Allows downloading a document by specifying its filename."
)
async def download_file(filename: str):
    """Allows downloading a document generated."""
    file_path = FILES_FOLDER / filename  
    if not file_path.is_file():
        raise HTTPException(status_code=404, detail="File not found")
    return FileResponse(file_path, filename=filename, headers={"Content-Disposition": "attachment"})




@app.post("/files/test/")
async def pdf2html_and_html2pdf(file: UploadFile = File(...)):
    """Recibe un PDF, extrae su diseño como HTML y lo reconstruye en un nuevo PDF."""
    
    # Definir rutas de archivos
    pdf_path = FILES_FOLDER / file.filename
    html_path = FILES_FOLDER / f"{file.filename}.html"
    pdf_output_path = FILES_FOLDER / f"reconstructed_{file.filename}"

    # Guardar el archivo PDF subido
    with open(pdf_path, "wb") as pdf_file:
        pdf_file.write(await file.read())

    # Convertir PDF a HTML utilizando pdf2htmlEX en Docker
    subprocess.run([
        "docker", "run", "--rm", "-v", f"{FILES_FOLDER}:/pdf", "pdf2htmlex/pdf2htmlex",
        "pdf2htmlEX", "--dest-dir", "/pdf", f"/pdf/{file.filename}"
    ], check=True)

    # Leer el contenido HTML generado
    with open(html_path, "r", encoding="utf-8") as html_file:
        source_html = html_file.read()

    # Convertir HTML a PDF utilizando xhtml2pdf
    with open(pdf_output_path, "wb") as output_pdf:
        pisa.CreatePDF(source_html, dest=output_pdf)

    return {
        "message": "Conversión exitosa",
        "pdf_url": file.filename,
        "html_url": html_path.name,
        "pdf2_url": pdf_output_path.name
    }

@app.delete(
    "/files/delete/{filename}",
    summary="Delete a file",
    description="Deletes a file from the server."
)
async def delete_file(filename: str):
    """Deletes a document by specifying its filename."""
    file_path = FILES_FOLDER / filename  

    if not file_path.is_file():
        raise HTTPException(status_code=404, detail="File not found")

    try:
        file_path.unlink()  # Usando pathlib para eliminar el archivo
        return {"message": f"File '{filename}' deleted successfully"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error deleting file: {str(e)}")

@app.get("/files/list",
    summary="List generated documents",
    description="Lists all the generated documents available for download.")
async def list_files():
    files = [
        {
            "filename": file.name,
        }
        for file in FILES_FOLDER.iterdir() if file.is_file()
    ]
    return {"files": files} if files else {"message": "No files stored.", "files": []}

@app.get("/files/text/{filename}", 
         response_model=FileContentResponse, 
         summary="Retrieve extracted file text", 
         description="Returns the plain text of a previously uploaded file.")
async def get_file_text(filename: str):
    file_path = FILES_FOLDER / filename
    if file_path.exists() and file_path.is_file():
        file_text = read_file_text(file_path)
        return FileContentResponse(filename=filename, file_content=file_text)
    
    raise HTTPException(status_code=404, detail="File not found.") 

if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=7121)
